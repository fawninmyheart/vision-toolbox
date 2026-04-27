#!/usr/bin/env python3
r"""Universal file analyzer — images, PDFs, videos, Office docs, code, and more.

Sends files to any OpenAI-compatible multimodal API for analysis.

Usage:
  vision.py image.jpg                              # NVIDIA free API (default)
  vision.py image.jpg --provider openai            # OpenAI
  vision.py image.jpg --provider ollama --model llama3.2-vision  # Local
  vision.py image.jpg --base-url http://x:8080/v1/chat/completions  # Custom
  vision.py --presets                              # List available presets

Env vars: NVIDIA_API_KEY, OPENAI_API_KEY, ANTHROPIC_API_KEY, or API_KEY
"""

import argparse, base64, io, json, mimetypes, os, sys, tempfile, subprocess
from pathlib import Path
from urllib.request import Request, urlopen
from urllib.error import URLError

# ── provider presets ─────────────────────────────────────────────────

PRESETS = {
    "nvidia": {
        "base_url": "https://integrate.api.nvidia.com/v1/chat/completions",
        "model": "moonshotai/kimi-k2.5",
        "key_env": "NVIDIA_API_KEY",
        "desc": "NVIDIA NIM free tier (Kimi K2.5, Qwen3.5 VL, Llama Vision…)",
    },
    "openai": {
        "base_url": "https://api.openai.com/v1/chat/completions",
        "model": "gpt-4o",
        "key_env": "OPENAI_API_KEY",
        "desc": "OpenAI (GPT-4o, GPT-4.1…)",
    },
    "anthropic": {
        "base_url": "https://api.anthropic.com/v1/messages",
        "model": "claude-sonnet-4-6",
        "key_env": "ANTHROPIC_API_KEY",
        "desc": "Anthropic (not OpenAI-compatible; use a proxy)",
    },
    "ollama": {
        "base_url": "http://localhost:11434/v1/chat/completions",
        "model": "llava:latest",
        "key_env": None,
        "desc": "Local Ollama (llava, llama3.2-vision, minicpm-v…)",
    },
    "lmstudio": {
        "base_url": "http://localhost:1234/v1/chat/completions",
        "model": "auto",
        "key_env": None,
        "desc": "Local LM Studio",
    },
    "openrouter": {
        "base_url": "https://openrouter.ai/api/v1/chat/completions",
        "model": "openai/gpt-4o",
        "key_env": "OPENROUTER_API_KEY",
        "desc": "OpenRouter (multi-provider gateway)",
    },
}

DEFAULT_PRESET = "nvidia"
MAX_TEXT_CHARS = 80000
PDF_MAX_PAGES = 5
PDF_DPI = 150

DEFAULT_PROMPT = "请详细分析这封文件的内容。对于图片：描述物体、人物、场景和文字。对于视频帧序列：描述场景变化和关键动作。对于文档：提取并总结关键信息。对于代码/数据：解释结构和逻辑。"

# ── file classifiers ─────────────────────────────────────────────────

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif"}
TEXT_EXTS = {
    ".txt", ".csv", ".tsv", ".json", ".xml", ".html", ".htm", ".md", ".markdown",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".scss", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".conf", ".log", ".sh", ".bash", ".zsh",
    ".rs", ".go", ".java", ".c", ".cpp", ".h", ".hpp", ".rb", ".php",
    ".sql", ".r", ".m", ".swift", ".kt", ".scala", ".lua", ".vim",
    ".dockerfile", ".makefile", ".cmake", ".tex", ".rst",
}
OFFICE_EXTS = {".docx", ".xlsx", ".pptx", ".ods", ".odt", ".odp", ".rtf"}
VIDEO_EXTS = {".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv", ".m4v"}
VIDEO_MAX_FRAMES = 8
VIDEO_INTERVAL = 5  # seconds between frames


def classify(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext in VIDEO_EXTS:
        return "video"
    if ext in IMAGE_EXTS:
        return "image"
    if ext == ".pdf":
        return "pdf"
    if ext in OFFICE_EXTS:
        return "office"
    if ext in TEXT_EXTS:
        return "text"
    # fallback: try reading as text
    try:
        with open(path, "r", encoding="utf-8") as f:
            f.read(1024)
        return "text"
    except Exception:
        pass
    # check mime
    mime, _ = mimetypes.guess_type(path)
    if mime and mime.startswith("image/"):
        return "image"
    if mime and mime.startswith("text/"):
        return "text"
    return "binary"


# ── converters ───────────────────────────────────────────────────────

def encode_image(path: str) -> str:
    mime, _ = mimetypes.guess_type(path)
    if not mime or not mime.startswith("image/"):
        mime = "image/jpeg"
    with open(path, "rb") as f:
        data = base64.b64encode(f.read()).decode()
    return f"data:{mime};base64,{data}"


def video_to_frames(video_path: str) -> list[str]:
    """Extract key frames from video, return list of JPEG paths."""
    tmpdir = tempfile.mkdtemp(prefix="vision-video-")
    # Get duration
    dur_result = subprocess.run(
        ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
         "-of", "csv=p=0", video_path],
        capture_output=True, text=True, timeout=15
    )
    duration = float(dur_result.stdout.strip() or 0)
    if duration <= 0:
        raise RuntimeError(f"Cannot determine video duration: {video_path}")

    # Calculate frame count
    frame_count = min(VIDEO_MAX_FRAMES, max(3, int(duration / VIDEO_INTERVAL)))
    interval = duration / (frame_count + 1)

    paths = []
    for i in range(frame_count):
        t = interval * (i + 1)
        out_path = f"{tmpdir}/frame_{i+1:02d}_{int(t)}s.jpg"
        result = subprocess.run(
            ["ffmpeg", "-ss", str(t), "-i", video_path,
             "-vframes", "1", "-q:v", "3", "-y", out_path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            paths.append(out_path)

    if not paths:
        raise RuntimeError(f"Failed to extract any frames from {video_path}")
    return paths


def pdf_to_images(pdf_path: str) -> list[str]:
    tmpdir = tempfile.mkdtemp(prefix="vision-")
    code = f"""
import fitz
doc = fitz.open('{pdf_path}')
for i in range(min(len(doc), {PDF_MAX_PAGES})):
    pix = doc[i].get_pixmap(dpi={PDF_DPI})
    pix.save('{tmpdir}/page_' + str(i+1) + '.jpg')
"""
    result = subprocess.run(
        ["micromamba", "run", "-n", "pdf-toolbox", "python3", "-c", code],
        capture_output=True, text=True, timeout=60
    )
    if result.returncode != 0:
        raise RuntimeError(f"PDF render failed: {result.stderr}")
    import glob
    return sorted(glob.glob(f"{tmpdir}/page_*.jpg"))


def extract_office_text(path: str) -> str:
    ext = Path(path).suffix.lower()
    code = ""
    if ext == ".docx":
        code = f"from docx import Document; doc=Document('{path}'); [print(p.text) for p in doc.paragraphs]"
    elif ext == ".xlsx":
        code = f"""
from openpyxl import load_workbook
wb = load_workbook('{path}', data_only=True)
for name in wb.sheetnames:
    ws = wb[name]
    print(f'--- Sheet: {name} ---')
    for row in ws.iter_rows(values_only=True):
        vals = [str(v) if v is not None else '' for v in row]
        if any(v for v in vals):
            print('\t'.join(vals))
"""
    elif ext == ".pptx":
        code = f"""
from pptx import Presentation
prs = Presentation('{path}')
for i, slide in enumerate(prs.slides):
    print(f'--- Slide {i+1} ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(text)
"""
    elif ext in (".odt", ".ods", ".odp", ".rtf"):
        return f"[需要 LibreOffice 转换的文档格式: {ext}，请安装 libreoffice]"
    else:
        return ""

    if not code:
        return ""

    result = subprocess.run(
        ["micromamba", "run", "-n", "pdf-toolbox", "python3", "-c", code],
        capture_output=True, text=True, timeout=30
    )
    if result.returncode != 0:
        raise RuntimeError(f"Office extraction failed: {result.stderr}")
    return result.stdout


def read_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read(MAX_TEXT_CHARS + 1)
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n\n... [truncated]"
    return text


# ── API ──────────────────────────────────────────────────────────────

def build_message(prompt: str, items: list[dict]) -> dict:
    content = [{"type": "text", "text": prompt}]
    for item in items:
        if item["kind"] == "image":
            content.append({
                "type": "image_url",
                "image_url": {"url": item["data"]}
            })
        elif item["kind"] == "text":
            content.append({
                "type": "text",
                "text": f"--- FILE: {item['name']} ---\n{item['data']}"
            })
    return {"role": "user", "content": content}


def call_api(api_key: str, base_url: str, model: str, messages: list[dict],
             max_tokens: int = 4096) -> dict:
    body = json.dumps({
        "model": model,
        "messages": messages,
        "max_tokens": max_tokens,
        "temperature": 0.2,
    }).encode()

    req = Request(base_url, data=body, headers={
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    })

    try:
        with urlopen(req, timeout=180) as resp:
            return json.loads(resp.read())
    except URLError as e:
        if hasattr(e, "read"):
            return json.loads(e.read())
        return {"error": str(e)}


# ── main ─────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Universal file analyzer via any OpenAI-compatible API")
    parser.add_argument("files", nargs="*", help="File paths to analyze")
    parser.add_argument("--prompt", default=DEFAULT_PROMPT, help="Custom prompt")
    parser.add_argument("--provider", default=DEFAULT_PRESET,
                        help=f"Provider preset (default: {DEFAULT_PRESET})")
    parser.add_argument("--base-url", help="Override API base URL")
    parser.add_argument("--model", help="Override model ID")
    parser.add_argument("--key", default=os.environ.get("API_KEY"),
                        help="API key (env: API_KEY, or provider-specific)")
    parser.add_argument("--max-tokens", type=int, default=4096)
    parser.add_argument("--presets", action="store_true", help="List provider presets")
    args = parser.parse_args()

    if args.presets:
        print("Available presets (--provider):\n")
        for name, cfg in PRESETS.items():
            marker = " [default]" if name == DEFAULT_PRESET else ""
            print(f"  {name}{marker}")
            print(f"    {cfg['desc']}")
            print(f"    model: {cfg['model']}")
            if cfg["key_env"]:
                print(f"    key:   ${cfg['key_env']}")
            print()
        return

    # Resolve provider
    preset = PRESETS.get(args.provider, PRESETS[DEFAULT_PRESET])
    base_url = args.base_url or preset["base_url"]
    model = args.model or preset["model"]
    api_key = args.key
    if not api_key and preset["key_env"]:
        api_key = os.environ.get(preset["key_env"])

    if not api_key:
        key_hint = preset["key_env"] or "API_KEY"
        print(f"ERROR: No API key. Set ${key_hint} or pass --key", file=sys.stderr)
        sys.exit(1)

    if not args.files:
        parser.error("at least one file is required")

    # Process files
    image_items = []
    text_parts = []
    stats = []

    for f in args.files:
        kind = classify(f)
        name = os.path.basename(f)
        try:
            if kind == "image":
                data = encode_image(f)
                image_items.append({"kind": "image", "name": name, "data": data})
                stats.append(f"IMG {name}")
            elif kind == "pdf":
                imgs = pdf_to_images(f)
                for img_path in imgs:
                    data = encode_image(img_path)
                    image_items.append({"kind": "image", "name": Path(img_path).name, "data": data})
                stats.append(f"PDF {name} ({len(imgs)} pages)")
            elif kind == "video":
                frames = video_to_frames(f)
                for img_path in frames:
                    data = encode_image(img_path)
                    image_items.append({"kind": "image", "name": Path(img_path).name, "data": data})
                stats.append(f"VID {name} ({len(frames)} frames)")
            elif kind == "office":
                extracted = extract_office_text(f)
                if extracted.startswith("[需要"):
                    image_items.append({"kind": "text", "name": name, "data": extracted})
                else:
                    text_parts.append(f"--- FILE: {name} ---\n{extracted}")
                stats.append(f"DOC {name} ({len(extracted)} chars)")
            elif kind == "text":
                text = read_text_file(f)
                text_parts.append(f"--- FILE: {name} ---\n{text}")
                stats.append(f"TXT {name} ({len(text)} chars)")
            else:
                stats.append(f"SKIP {name} (binary)")
                print(f"WARNING: Skipping binary file: {name}", file=sys.stderr)
        except Exception as e:
            print(f"ERROR processing {name}: {e}", file=sys.stderr)
            sys.exit(1)

    if not image_items and not text_parts:
        print("No processable files found.", file=sys.stderr)
        sys.exit(1)

    # Build prompt with text content prepended
    full_prompt = args.prompt
    if text_parts:
        full_prompt = "\n\n".join(text_parts) + "\n\n---\n" + full_prompt

    messages = [build_message(full_prompt, image_items)]
    print(f"[Sending] {' '.join(stats)}", file=sys.stderr)

    result = call_api(api_key, base_url, model, messages, args.max_tokens)

    if "error" in result:
        print(f"API Error: {result['error']}", file=sys.stderr)
        sys.exit(1)

    if "choices" in result:
        msg = result["choices"][0]["message"]
        content = msg.get("content", "")
        print(content)
    else:
        print(json.dumps(result, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
